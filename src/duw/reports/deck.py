"""Client deck.

An optional python-pptx summary deck: a title slide plus a handful of bullet
slides covering the counterparty, exposure, and the illustrative recommendation.
Charts are embedded as PNGs when kaleido is available. Pure Python; no Qt.
"""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from duw.domain.results import AnalysisResults
from duw.reports.interpreter import DISCLAIMER, interpret_exposure, recommend
from duw.reports.memo import (
    _counterparty_rows,
    _exposure_rows,
    _figure_png,
    _trade_rows,
)


def _bullets(slide_body, lines: list[str]) -> None:
    tf = slide_body.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line


def write_memo_pptx(results: AnalysisResults, path: str | Path) -> Path:
    """Write a client summary deck to ``path``."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    out = Path(path)
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]

    # Title slide.
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Counterparty Credit Underwriting"
    cp = results.counterparty
    slide.placeholders[1].text = (
        f"{cp.name} — educational portfolio model"
        if cp is not None
        else "Educational portfolio model"
    )

    # Trade + counterparty slide.
    slide = prs.slides.add_slide(bullet_layout)
    slide.shapes.title.text = "Trade & Counterparty"
    lines = [f"{k}: {v}" for k, v in _trade_rows(results)]
    lines += [f"{k}: {v}" for k, v in _counterparty_rows(results)]
    _bullets(slide.placeholders[1], lines)

    # Exposure slide with commentary and chart.
    slide = prs.slides.add_slide(bullet_layout)
    slide.shapes.title.text = "Exposure"
    _bullets(
        slide.placeholders[1],
        [f"{k}: {v}" for k, v in _exposure_rows(results)]
        + [interpret_exposure(results)],
    )
    from duw.ui.widgets.charts import exposure_figure

    png = _figure_png(exposure_figure(results.exposure))
    if png is not None:
        slide.shapes.add_picture(
            BytesIO(png), Inches(5.2), Inches(1.6), width=Inches(4.2)
        )

    # Recommendation slide.
    rec = recommend(results)
    slide = prs.slides.add_slide(bullet_layout)
    slide.shapes.title.text = "Recommendation"
    _bullets(slide.placeholders[1], [rec.verdict, rec.rationale])

    # Disclaimer slide.
    slide = prs.slides.add_slide(bullet_layout)
    slide.shapes.title.text = "Disclaimer"
    body = slide.placeholders[1]
    _bullets(body, [DISCLAIMER])
    for para in body.text_frame.paragraphs:
        para.font.size = Pt(11)

    prs.save(str(out))
    return out
